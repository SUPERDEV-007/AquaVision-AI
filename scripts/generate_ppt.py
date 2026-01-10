
import collections 
import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
import os

# Define colors
PRIMARY_COLOR = RGBColor(30, 136, 229)    # Blue #1E88E5
SECONDARY_COLOR = RGBColor(67, 160, 71)   # Green #43A047
ACCENT_COLOR = RGBColor(251, 140, 0)      # Orange #FB8C00
TEXT_COLOR = RGBColor(33, 33, 33)         # Dark Gray #212121
BG_COLOR = RGBColor(245, 245, 245)        # Light Gray #F5F5F5
WHITE = RGBColor(255, 255, 255)

def create_presentation():
    prs = Presentation()
    
    # Set slide dimensions to widescreen 16:9
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    def add_slide(layout_index=1, title="", content=""):
        slide_layout = prs.slide_layouts[layout_index]
        slide = prs.slides.add_slide(slide_layout)
        
        # Set background
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = BG_COLOR

        if slide.shapes.title:
            slide.shapes.title.text = title
            # Title styling
            title_tf = slide.shapes.title.text_frame
            title_p = title_tf.paragraphs[0]
            title_p.font.name = 'Arial'
            title_p.font.size = Pt(40)
            title_p.font.bold = True
            title_p.font.color.rgb = PRIMARY_COLOR

        if content:
            # Depending on layout, content might be in placeholders[1]
            if len(slide.placeholders) > 1:
                body = slide.placeholders[1]
                body.text = content
                
                # Content styling
                for paragraph in body.text_frame.paragraphs:
                    paragraph.font.name = 'Arial'
                    paragraph.font.size = Pt(24)
                    paragraph.font.color.rgb = TEXT_COLOR
                    paragraph.space_after = Pt(14)

        return slide

    # --- SLIDE 1: Title Slide ---
    # Using Layout 0 (Title Slide)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    background = slide1.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR # Dark Blue background for title
    
    title = slide1.shapes.title
    subtitle = slide1.placeholders[1]
    
    title.text = "AI-Powered Fish Health System"
    subtitle.text = "Aquaculture Disease Detection Using Deep Learning\n\n- Real-time disease classification\n- 7-class disease detection\n- Explainable AI & Edge Deployment"
    
    # Style Title Slide
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.size = Pt(54)
    title.text_frame.paragraphs[0].font.bold = True
    
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = WHITE
        p.font.size = Pt(28)

    # --- SLIDE 2: Problem Statement ---
    s2 = add_slide(1, "Problem & Solution")
    tf = s2.placeholders[1].text_frame
    tf.clear()
    
    p = tf.add_paragraph()
    p.text = "Challenges in Aquaculture:"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = ACCENT_COLOR
    
    bullets = [
        "Manual disease detection is time-consuming and prone to errors",
        "Early detection is critical but expert knowledge is scarce",
        "Real-time monitoring needed for large farms"
    ]
    for b in bullets:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1
        p.font.size = Pt(24)

    p = tf.add_paragraph()
    p.text = "\nOur Solution:"
    p.font.bold = True
    p.font.size = Pt(28)
    p.font.color.rgb = SECONDARY_COLOR
    
    bullets2 = [
        "AI-powered automated detection using Computer Vision",
        "Instant diagnosis with high accuracy (76.55%)",
        "Explainable predictions (Grad-CAM) & Offline Edge Deployment"
    ]
    for b in bullets2:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1
        p.font.size = Pt(24)

    # --- SLIDE 3: Project Overview ---
    s3 = add_slide(1, "Project Overview")
    tf = s3.placeholders[1].text_frame
    tf.text = "Core Objective:\nDevelop an explainable, real-time AI system for multi-class image classification of freshwater fish diseases."
    
    p = tf.add_paragraph()
    p.text = "\nKey Features:"
    p.font.bold = True
    
    features = [
        "✅ 7 Disease Classes: Healthy + 6 common diseases",
        "✅ Architecture: MobileNetV2 (Transfer Learning)",
        "✅ High Accuracy: 76.55% on test set",
        "✅ Explainable AI: Grad-CAM visualizations",
        "✅ Edge Deployment: TFLite (2.40 MB)"
    ]
    for f in features:
        p = tf.add_paragraph()
        p.text = f
        p.level = 1

    # --- SLIDE 4: Disease Classes ---
    s4 = add_slide(1, "Disease Classes (7 Categories)")
    # We can create a table here or just list them. A table is nicer but harder to code blindly.
    # We'll use a text list for safety and simplicity in this script.
    tf = s4.placeholders[1].text_frame
    tf.text = "Dataset: 2,400+ images, balanced distribution"
    
    classes = [
        "🟢 Healthy Fish (Health Status)",
        "🔴 Aeromoniasis (Bacterial)",
        "🔴 Gill Disease (Bacterial)",
        "🔴 Red Disease (Bacterial)",
        "🟡 Saprolegniasis (Fungal)",
        "🟠 Parasitic Disease (Parasitic)",
        "🔵 White Tail Disease (Viral)"
    ]
    for c in classes:
        p = tf.add_paragraph()
        p.text = c
        p.level = 1
        p.space_after = Pt(10)

    # --- SLIDE 5: Technical Architecture ---
    s5 = add_slide(1, "Technical Architecture")
    tf = s5.placeholders[1].text_frame
    tf.text = "Deep Learning Pipeline:"
    
    pipeline = [
        "Input Image (224×224×3)",
        "⬇ MobileNetV2 (Pre-trained on ImageNet)",
        "⬇ Feature Extraction (Frozen layers)",
        "⬇ Global Average Pooling & Dropout (0.3)",
        "⬇ Dense Layer (7 classes) -> Softmax",
        "⬇ Disease Prediction"
    ]
    for item in pipeline:
        p = tf.add_paragraph()
        p.text = item
        p.level = 1
        p.font.size = Pt(22)

    p = tf.add_paragraph()
    p.text = "\nStack: TensorFlow/Keras 3, Adam Optimizer, Grad-CAM"
    p.font.bold = True

    # --- SLIDE 6 & 7: Performance ---
    s6 = add_slide(1, "Model Performance")
    tf = s6.placeholders[1].text_frame
    tf.text = "Metrics (Test Set - 469 images):"
    
    metrics = [
        "🏆 Accuracy: 76.55%",
        "🎯 Precision: 88.42%",
        "🔍 Recall: 58.64%",
        "⚖️ F1-Score: 76.41%"
    ]
    for m in metrics:
        p = tf.add_paragraph()
        p.text = m
        p.level = 1
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = PRIMARY_COLOR

    p = tf.add_paragraph()
    p.text = "\nBest Performing Classes:"
    p.font.size = Pt(24)
    p.font.color.rgb = TEXT_COLOR
    
    best = [
        "⭐ Healthy: 93% Recall (Critical for safety)",
        "⭐ Gill Disease: 88% Precision"
    ]
    for b in best:
        p = tf.add_paragraph()
        p.text = b
        p.level = 1

    # --- SLIDE 8: Explainable AI (Grad-CAM) ---
    s8 = add_slide(5, "Explainable AI - Grad-CAM") # Layout 5 is usually Title only
    # Add image if exists
    img_path = os.path.join("results", "gradcam_visualization.png")
    if os.path.exists(img_path):
        s8.shapes.add_picture(img_path, Inches(1), Inches(2), width=Inches(11))
    else:
        # Fallback text
        txBox = s8.shapes.add_textbox(Inches(1), Inches(2), Inches(10), Inches(4))
        tf = txBox.text_frame
        tf.text = "Grad-CAM (Gradient-weighted Class Activation Mapping) highlights image regions influencing predictions.\n(Image not found in results folder)"

    # --- SLIDE 9: Edge Deployment ---
    s9 = add_slide(1, "Edge-AI Deployment")
    tf = s9.placeholders[1].text_frame
    tf.text = "TensorFlow Lite Conversion"
    
    stats = [
        "Original Model: ~8.65 MB",
        "⬇ Optimized TFLite: 2.40 MB (Quantized)",
        "⬇ 72% Size Reduction"
    ]
    for s in stats:
        p = tf.add_paragraph()
        p.text = s
        p.level = 1
        p.font.bold = True
    
    p = tf.add_paragraph()
    p.text = "\ncapabilities:"
    p.font.bold = True
    
    caps = [
        "✅ Offline Operation (No internet)",
        "✅ Real-time Inference on Mobile/Edge devices",
        "✅ Low Resource Consumption"
    ]
    for c in caps:
        p = tf.add_paragraph()
        p.text = c
        p.level = 1

    # --- SLIDE 11: Achievements ---
    s11 = add_slide(1, "Key Achievements")
    tf = s11.placeholders[1].text_frame
    tf.text = "System Status: Fully Operational"
    
    achievements = [
        "✅ Balanced stratified dataset (7 classes)",
        "✅ Keras 3 Compatibility & Modern UI",
        "✅ Real-time Grad-CAM Heatmaps",
        "✅ Edge-Ready TFLite Model (2.4MB)",
        "✅ 76.55% Accuracy verified"
    ]
    for a in achievements:
        p = tf.add_paragraph()
        p.text = a
        p.level = 1
        p.space_after = Pt(20)

    # --- SLIDE 12: Use Cases ---
    s12 = add_slide(1, "Real-World Applications")
    tf = s12.placeholders[1].text_frame
    
    uses = [
        "🌊 Aquaculture Farms: Daily monitoring & disease prevention",
        "🔬 Research: Disease pattern analysis",
        "📱 Mobile Diagnosis: Farmers with smartphones (Offline)",
        "🐟 Hatcheries: Quality control & breeding selection"
    ]
    for u in uses:
        p = tf.add_paragraph()
        p.text = u
        p.space_after = Pt(20)
        p.font.size = Pt(26)

    # --- SLIDE 17: Conclusion ---
    s17 = add_slide(1, "Conclusion")
    tf = s17.placeholders[1].text_frame
    tf.text = "Impact & Summary"
    
    conclusions = [
        "We built a fully functional, end-to-end AI system.",
        "Solves a critical problem in aquaculture with 76% accuracy.",
        "Transparent decision making (Explainable AI).",
        "Ready for real-world use on edge devices."
    ]
    for c in conclusions:
        p = tf.add_paragraph()
        p.text = c
        p.level = 1
        p.font.size = Pt(28)
        p.space_after = Pt(24)

    # --- SLIDE 18: Thank You ---
    slide_end = prs.slides.add_slide(prs.slide_layouts[0])
    background = slide_end.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR
    
    title = slide_end.shapes.title
    title.text = "Thank You!"
    title.text_frame.paragraphs[0].font.color.rgb = WHITE
    title.text_frame.paragraphs[0].font.size = Pt(70)
    
    subtitle = slide_end.placeholders[1]
    subtitle.text = "Questions & Discussion\n\nFish Disease Detection System"
    for p in subtitle.text_frame.paragraphs:
        p.font.color.rgb = WHITE

    # Save
    prs.save('Fish_Disease_Detection_Presentation.pptx')
    print("Presentation saved successfully.")

if __name__ == "__main__":
    create_presentation()
